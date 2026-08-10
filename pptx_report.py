#!/usr/bin/env python3
"""
Gera a apresentação PowerPoint de diagnóstico financeiro (python-pptx puro —
sem depender de Node/pptxgenjs, para o app rodar só com Python) a partir do
resultado de pipeline.consolidate() + pipeline.compute_indicators(). Mesma
paleta "Midnight Executive" usada no deck original do Grupo CLA/Claudio
(navy + azul-gelo), sem qualquer marca institucional — pronta para
apresentar a bancos sob a identidade de quem gerar o relatório.

Assim como no docx_report.py, os números são 100% automáticos e os "pontos
de atenção" seguem regras de referência simples (THRESHOLDS, compartilhado
com docx_report). Para uma narrativa mais fina de um caso específico, use a
skill "diagnostico-bancario-grupo-cla-claudio" com o Claude.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

import docx_report as dr  # reaproveita fmt_mi / fmt_pct / fmt_x / THRESHOLDS

NAVY = RGBColor(0x1E, 0x27, 0x61)
NAVY_DARK = RGBColor(0x14, 0x1B, 0x45)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
ICE_LIGHT = RGBColor(0xEE, 0xF3, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x5B, 0x64, 0x72)
POS = RGBColor(0x2C, 0x5F, 0x2D)
POS_BG = RGBColor(0xE4, 0xEE, 0xE4)
CRIT = RGBColor(0x8A, 0x33, 0x24)
CRIT_BG = RGBColor(0xF4, 0xE3, 0xDE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def add_text(slide, x, y, w, h, text, size=14, bold=False, italic=False, color=RGBColor(0x22, 0x22, 0x22),
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_rect(slide, x, y, w, h, fill=ICE_LIGHT, line=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def section_title(slide, title, subtitle=None):
    add_text(slide, Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.7), title,
             size=28, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.4), subtitle,
                  size=13, italic=True, color=GRAY)


def footer(slide, extra=""):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(9.5), Inches(0.3), extra,
              size=9, italic=True, color=GRAY)
    add_text(slide, Inches(0.5), Inches(7.1), Inches(12.33), Inches(0.3), "Diagnóstico Financeiro Consolidado",
              size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def stat_card(slide, x, y, w, h, value, label, value_color=NAVY):
    add_rect(slide, x, y, w, h, fill=ICE_LIGHT)
    add_text(slide, x + Inches(0.1), y + Inches(0.12), w - Inches(0.2), h * 0.55, value,
              size=24, bold=True, color=value_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
    add_text(slide, x + Inches(0.1), y + h * 0.6, w - Inches(0.2), h * 0.35, label,
              size=11.5, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


def add_indicator_table(slide, x, y, w, header, rows, col_widths=None, font_size=11):
    n_rows = len(rows) + 1
    n_cols = len(header)
    h = Inches(0.4) * n_rows
    gtable = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = gtable.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for j, lbl in enumerate(header):
        cell = table.cell(0, j)
        cell.text = lbl
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = WHITE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ICE_LIGHT if i % 2 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(font_size)
            run.font.bold = (j == 0)
            run.font.color.rgb = NAVY if j == 0 else RGBColor(0x22, 0x22, 0x22)
    return gtable


def add_full_statement_table(slide, x, y, w, h, row_template, values_by_period, labels, font_size=8):
    """Tabela completa de um demonstrativo (BP/DRE/DFC), uma linha por conta,
    uma coluna por período + Agregado. Usada nos slides de anexo (relatórios
    completos), separado dos slides de indicadores-resumo."""
    n_rows = len(row_template) + 1
    n_cols = 1 + len(labels)
    gtable = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = gtable.table
    table.first_row = False
    label_w = Emu(int(w * 0.42))
    table.columns[0].width = label_w
    data_col_w = Emu(int((w - label_w) / max(len(labels), 1)))
    for j in range(len(labels)):
        table.columns[1 + j].width = data_col_w

    row_h = Emu(int(h / n_rows))
    for r in range(n_rows):
        table.rows[r].height = row_h

    # header
    hcell = table.cell(0, 0)
    hcell.text = "Descrição"
    hcell.fill.solid(); hcell.fill.fore_color.rgb = NAVY
    for j, lbl in enumerate(labels):
        cell = table.cell(0, 1 + j)
        cell.text = lbl
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for j in range(n_cols):
        cell = table.cell(0, j)
        cell.margin_left = Emu(18000); cell.margin_right = Emu(18000)
        cell.margin_top = Emu(4000); cell.margin_bottom = Emu(4000)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = WHITE

    for i, (label, depth, key) in enumerate(row_template, start=1):
        desc_cell = table.cell(i, 0)
        desc_cell.text = ("    " * depth) + label
        if key is None:
            fill = ICE
        elif depth == 0:
            fill = ICE_LIGHT
        else:
            fill = WHITE
        for j in range(n_cols):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            cell.margin_left = Emu(18000); cell.margin_right = Emu(18000)
            cell.margin_top = Emu(2000); cell.margin_bottom = Emu(2000)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = desc_cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(font_size)
        run.font.bold = (depth == 0)
        run.font.color.rgb = NAVY if depth == 0 else RGBColor(0x33, 0x33, 0x33)

        if key is None:
            for j in range(1, n_cols):
                table.cell(i, j).text = ""
            continue
        for j in range(len(labels)):
            val = values_by_period[j].get(key, 0.0) / 1_000_000.0
            cell = table.cell(i, 1 + j)
            cell.text = dr.fmt_mi(val)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(font_size)
            run.font.bold = (depth == 0)
            run.font.color.rgb = NAVY if depth == 0 else RGBColor(0x33, 0x33, 0x33)
    return gtable


def add_statement_slides(prs, title, subtitle, row_template, values_by_period, labels, note=None):
    """Divide o demonstrativo em quantos slides forem necessários para caber
    (cerca de 16 linhas por slide, em fonte pequena, para não estourar)."""
    CHUNK = 16
    chunks = [row_template[k:k + CHUNK] for k in range(0, len(row_template), CHUNK)]
    total = len(chunks)
    for idx, chunk in enumerate(chunks, start=1):
        s = blank_slide(prs)
        sub = subtitle if total == 1 else f"{subtitle} (parte {idx}/{total})"
        section_title(s, title, sub)
        add_full_statement_table(s, Inches(0.5), Inches(1.55), Inches(12.33), Inches(5.5),
                                   chunk, values_by_period, labels, font_size=8.5)
        if note and idx == total:
            add_text(s, Inches(0.5), Inches(7.12), Inches(12.33), Inches(0.3), note,
                      size=8, italic=True, color=GRAY)
        else:
            footer(s)


COMMENT_FILL = RGBColor(0xF3, 0xF0, 0xE6)


def add_dre_variacao_slides(prs, dre_variacao_pares, dre_variacao_comentarios, labels):
    """Tabela com a variação de cada conta da DRE entre cada par de períodos
    consecutivos (período a período — não só primeiro x último), em R$
    milhões e %. Toda linha marcada como 'relevante' (>=20% ou >=R$0,5mi,
    calculado em build_dre_variacao_pares) que tiver um comentário do usuário
    ganha uma linha extra mesclada abaixo, com esse comentário. Só entra no
    PPT se o usuário marcar a opção no app ("Deseja levar a variação de
    resultado para o diagnóstico?")."""
    if not dre_variacao_pares:
        return
    dre_variacao_comentarios = dre_variacao_comentarios or {}

    CHUNK = 12  # em itens de render (linha de dado ou linha de comentário)
    for par in dre_variacao_pares:
        idx_a, idx_b = par["idx_a"], par["idx_b"]
        label_a, label_b = par["label_a"], par["label_b"]

        # monta a lista de itens a renderizar: cada linha de dado, seguida
        # (se houver) da linha de comentário do usuário para aquela conta.
        items = []
        for linha in par["linhas"]:
            items.append(("data", linha))
            if linha["relevante"]:
                comentario = dre_variacao_comentarios.get(f"{par['idx_a']}|{linha['key']}", "").strip()
                if comentario:
                    items.append(("comment", comentario))

        chunks = [items[k:k + CHUNK] for k in range(0, len(items), CHUNK)]
        total = len(chunks)
        for idx, chunk in enumerate(chunks, start=1):
            s = blank_slide(prs)
            sub = f"{label_a}  →  {label_b}"
            if total > 1:
                sub += f" (parte {idx}/{total})"
            section_title(s, "Variação das Contas de Resultado (DRE)", sub)

            x, y, w, h = Inches(0.5), Inches(1.55), Inches(12.33), Inches(5.5)
            n_rows = len(chunk) + 1
            n_cols = 5
            gtable = s.shapes.add_table(n_rows, n_cols, x, y, w, h)
            table = gtable.table
            col_widths = [Emu(int(w * 0.40)), Emu(int(w * 0.15)), Emu(int(w * 0.15)), Emu(int(w * 0.15)), Emu(int(w * 0.15))]
            for j, cw in enumerate(col_widths):
                table.columns[j].width = cw
            row_h = Emu(int(h / n_rows))
            for r in range(n_rows):
                table.rows[r].height = row_h

            header = ["Conta", label_a, label_b, "Variação (R$ mi)", "Variação (%)"]
            for j, lbl in enumerate(header):
                cell = table.cell(0, j)
                cell.text = lbl
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                cell.margin_left = Emu(18000); cell.margin_right = Emu(18000)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                run = p.runs[0]
                run.font.bold = True
                run.font.size = Pt(9)
                run.font.color.rgb = WHITE

            for i, (kind, payload) in enumerate(chunk, start=1):
                if kind == "comment":
                    table.cell(i, 0).merge(table.cell(i, n_cols - 1))
                    cell = table.cell(i, 0)
                    cell.text = f"Comentário: {payload}"
                    cell.fill.solid(); cell.fill.fore_color.rgb = COMMENT_FILL
                    cell.margin_left = Emu(36000); cell.margin_right = Emu(18000)
                    cell.margin_top = Emu(2000); cell.margin_bottom = Emu(2000)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.LEFT
                    run = p.runs[0] if p.runs else p.add_run()
                    run.font.size = Pt(8.5)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    continue

                linha = payload
                label, depth = linha["label"], linha["depth"]
                v0, v1, diff, pct = linha["v0"], linha["v1"], linha["diff"], linha["pct"]
                marca = " ⚠" if linha["relevante"] else ""
                vals = [("    " * depth) + label + marca, dr.fmt_mi(v0), dr.fmt_mi(v1), dr.fmt_mi(diff), dr.fmt_pct(pct)]
                fill = ICE_LIGHT if depth == 0 else WHITE
                for j, v in enumerate(vals):
                    cell = table.cell(i, j)
                    cell.text = str(v)
                    cell.fill.solid(); cell.fill.fore_color.rgb = fill
                    cell.margin_left = Emu(18000); cell.margin_right = Emu(18000)
                    cell.margin_top = Emu(2000); cell.margin_bottom = Emu(2000)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.CENTER
                    run = p.runs[0] if p.runs else p.add_run()
                    run.font.size = Pt(9)
                    run.font.bold = (depth == 0)
                    run.font.color.rgb = NAVY if depth == 0 else RGBColor(0x33, 0x33, 0x33)
            if idx == total:
                add_text(s, Inches(0.5), Inches(7.12), Inches(12.33), Inches(0.3),
                          "⚠ = variação acima de 20% ou R$ 0,5 mi entre os dois períodos. Comparação período a período (não é o Agregado).",
                          size=8, italic=True, color=GRAY)
            else:
                footer(s)



def build_presentation(statements, indicators, out_path, group_name="Empresa/Grupo",
                        purpose="Preparado para apresentação a parceiros bancários — finalidade: captação de recursos",
                        prepared_by="Consultec — Escritório de Contabilidade",
                        incluir_variacao_dre=False,
                        dre_variacao_pares=None,
                        dre_variacao_comentarios=None):
    prs = new_prs()
    labels = statements["period_labels"]
    n = len(labels)
    ind = indicators["por_periodo"]
    agg_idx = indicators["agg_idx"]
    last = ind[indicators["last_flow_idx"]]
    flow_labels = labels[:-1]

    # ---------------- SLIDE 1: CAPA ----------------
    s = blank_slide(prs, bg=NAVY_DARK)
    add_rect(s, Inches(5.87), Inches(0.85), Inches(1.6), Inches(1.6), fill=NAVY)
    add_text(s, Inches(5.87), Inches(0.85), Inches(1.6), Inches(1.6), "$", size=44, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.8), Inches(2.85), Inches(11.73), Inches(0.9), "DIAGNÓSTICO FINANCEIRO CONSOLIDADO",
              size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(3.7), Inches(11.73), Inches(0.6), group_name,
              size=22, bold=True, color=ICE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(4.4), Inches(11.73), Inches(0.4),
              f"Base: Balanço, DRE e DFC consolidados — {', '.join(flow_labels)}",
              size=14, italic=True, color=RGBColor(0xAF, 0xC1, 0xE8), align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(4.85), Inches(11.73), Inches(0.4), purpose,
              size=12.5, italic=True, color=RGBColor(0x8F, 0xA5, 0xD6), align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(6.7), Inches(11.73), Inches(0.35),
              "Consolidação por soma simples das empresas informadas, sem eliminação de saldos intercompanhias",
              size=10, italic=True, color=RGBColor(0x74, 0x88, 0xBD), align=PP_ALIGN.CENTER)

    # ---------------- SLIDE 2: SUMÁRIO EXECUTIVO ----------------
    s = blank_slide(prs)
    section_title(s, "Sumário Executivo", f"Números-chave — {', '.join(flow_labels)}")
    crescimento = indicators.get("crescimento_pct", {})
    resumo = (f"{group_name} encerrou o período coberto com receita líquida de {dr.fmt_mi(ind[agg_idx]['receita_liquida_mi'])} "
              f"e patrimônio líquido de {dr.fmt_mi(ind[agg_idx]['pl_total_mi'])}.")
    if crescimento:
        resumo += (f" Receita líquida variou {dr.fmt_pct(crescimento.get('receita_liquida', float('nan')))} "
                   f"e EBITDA variou {dr.fmt_pct(crescimento.get('ebitda', float('nan')))} entre o primeiro e o "
                   f"último período coberto.")
    add_text(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.8), resumo, size=13.5, color=RGBColor(0x33, 0x33, 0x33))

    cards = [
        (dr.fmt_mi(ind[agg_idx]["receita_liquida_mi"]), "Receita Líquida"),
        (dr.fmt_mi(ind[agg_idx]["ebitda_mi"]), "EBITDA"),
        (dr.fmt_pct(ind[agg_idx]["margem_ebitda_pct"]), "Margem EBITDA"),
        (dr.fmt_mi(ind[agg_idx]["resultado_liquido_mi"]), "Resultado Líquido"),
        (dr.fmt_mi(ind[agg_idx]["ativo_total_mi"]), "Ativo Total"),
        (dr.fmt_mi(ind[agg_idx]["pl_total_mi"]), "Patrimônio Líquido"),
    ]
    card_w, card_h, gap = Inches(3.85), Inches(1.55), Inches(0.25)
    start_x, start_y = Inches(0.6), Inches(2.55)
    for i, (val, lbl) in enumerate(cards):
        col, row = i % 3, i // 3
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)
        stat_card(s, x, y, card_w, card_h, val, lbl)
    footer(s)

    # ---------------- SLIDE 3: TRAJETÓRIA (gráfico) ----------------
    if crescimento:
        s = blank_slide(prs)
        section_title(s, "Trajetória entre Períodos", "Crescimento percentual — primeiro vs. último período coberto")
        chart_data = CategoryChartData()
        chart_data.categories = ["Receita Líquida", "EBITDA", "Resultado Líquido"]
        vals = [crescimento.get("receita_liquida", 0), crescimento.get("ebitda", 0), crescimento.get("resultado_liquido", 0)]
        chart_data.add_series("Crescimento (%)", [v if v == v else 0 for v in vals])
        gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.7), Inches(7.4), Inches(4.9), chart_data)
        chart = gframe.chart
        chart.has_legend = False
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            plot.data_labels.number_format = '0.0"%"'
            plot.data_labels.number_format_is_linked = False
            series = plot.series[0]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = NAVY
        except Exception:
            pass

        note_x, note_w = Inches(8.25), Inches(4.5)
        add_rect(s, note_x, Inches(1.7), note_w, Inches(4.9), fill=ICE_LIGHT)
        add_text(s, note_x + Inches(0.3), Inches(1.95), note_w - Inches(0.6), Inches(0.4), "Leitura", size=15, bold=True, color=NAVY)
        leitura = (f"Receita Líquida: {dr.fmt_pct(crescimento.get('receita_liquida', float('nan')))}\n"
                   f"EBITDA: {dr.fmt_pct(crescimento.get('ebitda', float('nan')))}\n"
                   f"Resultado Líquido: {dr.fmt_pct(crescimento.get('resultado_liquido', float('nan')))}")
        add_text(s, note_x + Inches(0.3), Inches(2.45), note_w - Inches(0.6), Inches(3.9), leitura, size=13, color=RGBColor(0x33, 0x33, 0x33))

    # ---------------- SLIDE 4: LIQUIDEZ E MARGENS ----------------
    s = blank_slide(prs)
    section_title(s, "Liquidez e Margens", "Indicadores por período coberto")
    header = ["Indicador"] + labels
    rows = [
        ["Liquidez Corrente"] + [dr.fmt_x(x["liquidez_corrente"]) for x in ind],
        ["Margem Bruta"] + [dr.fmt_pct(x["margem_bruta_pct"]) for x in ind],
        ["Margem EBITDA"] + [dr.fmt_pct(x["margem_ebitda_pct"]) for x in ind],
        ["Margem Líquida"] + [dr.fmt_pct(x["margem_liquida_pct"]) for x in ind],
    ]
    add_indicator_table(s, Inches(0.6), Inches(1.7), Inches(10.5), header, rows)
    footer(s)

    # ---------------- SLIDE 5: ENDIVIDAMENTO E ALAVANCAGEM ----------------
    s = blank_slide(prs)
    section_title(s, "Endividamento e Alavancagem", "Estrutura de capital e capacidade de pagamento")
    rows = [
        ["Endividamento Geral (Passivo/Ativo)"] + [dr.fmt_pct(x["endividamento_geral_pct"]) for x in ind],
        ["Patrimônio Líquido / Ativo Total"] + [dr.fmt_pct(x["pl_sobre_ativo_pct"]) for x in ind],
        ["Dívida Financeira Bruta"] + [dr.fmt_mi(x["divida_bruta_mi"]) for x in ind],
        ["Dívida Financeira Líquida"] + [dr.fmt_mi(x["divida_liquida_mi"]) for x in ind],
        ["Concentração de dívida no curto prazo"] + [dr.fmt_pct(x["concentracao_divida_cp_pct"]) for x in ind],
    ]
    add_indicator_table(s, Inches(0.6), Inches(1.7), Inches(7.6), header, rows)

    rx, rw = Inches(8.5), Inches(4.25)
    add_rect(s, rx, Inches(1.7), rw, Inches(4.95), fill=ICE_LIGHT)
    add_text(s, rx + Inches(0.25), Inches(2.0), rw - Inches(0.5), Inches(0.6), "Alavancagem\n(Dívida Líquida/EBITDA anualizado)",
              size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    alav_color = CRIT if (indicators["alavancagem_x"] == indicators["alavancagem_x"] and indicators["alavancagem_x"] > dr.THRESHOLDS["alavancagem_max_x"]) else POS
    add_text(s, rx + Inches(0.25), Inches(2.9), rw - Inches(0.5), Inches(1.0), dr.fmt_x(indicators["alavancagem_x"]),
              size=42, bold=True, color=alav_color, align=PP_ALIGN.CENTER)
    add_text(s, rx + Inches(0.25), Inches(4.0), rw - Inches(0.5), Inches(0.6),
              f"Fator de anualização usado: {indicators['fator_anualizacao']:.2f}x", size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, rx + Inches(0.25), Inches(4.7), rw - Inches(0.5), Inches(1.7),
              f"ROE anualizado: {dr.fmt_pct(indicators['roe_anualizado_pct'])}\nROA anualizado: {dr.fmt_pct(indicators['roa_anualizado_pct'])}\n"
              f"Referência de conforto usual dos bancos: até {dr.THRESHOLDS['alavancagem_max_x']:.1f}x.",
              size=12, color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.CENTER)
    footer(s)

    # ---------------- SLIDE 6: FLUXO DE CAIXA ----------------
    s = blank_slide(prs)
    section_title(s, "Fluxo de Caixa (DFC)", "Caixa gerado/consumido por atividade, método indireto")
    rows = [
        ["Caixa Operacional"] + [dr.fmt_mi(x["caixa_operacional_mi"]) for x in ind],
        ["Caixa de Investimento"] + [dr.fmt_mi(x["caixa_investimento_mi"]) for x in ind],
        ["Caixa de Financiamento"] + [dr.fmt_mi(x["caixa_financiamento_mi"]) for x in ind],
        ["Variação de Caixa no Período"] + [dr.fmt_mi(x["variacao_caixa_mi"]) for x in ind],
    ]
    add_indicator_table(s, Inches(0.6), Inches(1.7), Inches(11.5), header, rows)
    footer(s)

    # ---------------- VARIAÇÃO DAS CONTAS DE RESULTADO (opcional, a pedido do usuário no app) ----------------
    if incluir_variacao_dre:
        pares = dre_variacao_pares
        if pares is None:
            # fallback: se o app não passou os pares já calculados (ex.: chamada
            # direta desta função), recalcula período a período com os limiares padrão.
            import pipeline as _p
            pares = _p.build_dre_variacao_pares(statements, indicators, labels)
        add_dre_variacao_slides(prs, pares, dre_variacao_comentarios, labels)

    # ---------------- ANEXO: DEMONSTRATIVOS COMPLETOS (BP / DRE / DFC) ----------------
    # Sempre incluídos, com todas as linhas de cada demonstrativo (não só o resumo de
    # indicadores acima) — pedido explícito para que o PPT nunca saia sem os relatórios.
    add_statement_slides(
        prs, "Balanço Patrimonial Consolidado", "Anexo — todas as contas, por período (R$ milhões)",
        statements["bp_rows"], statements["bp"], labels,
        note="Balanço é uma fotografia num instante — a coluna \"Agregado\" repete a posição do último período, não soma.")
    add_statement_slides(
        prs, "Demonstração do Resultado (DRE) Consolidada", "Anexo — todas as contas, por período (R$ milhões)",
        statements["dre_rows"], statements["dre"], labels,
        note="\"Agregado\" soma todos os períodos incluídos (DRE é demonstração de fluxo).")
    add_statement_slides(
        prs, "Demonstração dos Fluxos de Caixa (DFC) Consolidada", "Anexo — método indireto, por período (R$ milhões)",
        statements["dfc_rows"], statements["dfc"], labels,
        note="\"Agregado\" soma os fluxos de todos os períodos; Caixa Inicial/Final não são somados (são saldos).")

    # ---------------- SLIDE 7: PONTOS DE ATENÇÃO ----------------
    s = blank_slide(prs)
    section_title(s, "Pontos de Atenção", "Sinalização automática por regras de referência de mercado")
    alertas = []
    if last["liquidez_corrente"] < dr.THRESHOLDS["liquidez_corrente_min"]:
        alertas.append(("Liquidez Corrente abaixo de 1,0x",
                         f"{dr.fmt_x(last['liquidez_corrente'])} no último período — ativo circulante não cobre o passivo circulante."))
    if last["cobertura_juros_x"] < dr.THRESHOLDS["cobertura_juros_min_x"]:
        alertas.append(("Cobertura de juros apertada",
                         f"EBITDA cobre {dr.fmt_x(last['cobertura_juros_x'])} das despesas financeiras — referência usual: acima de {dr.THRESHOLDS['cobertura_juros_min_x']:.1f}x."))
    if indicators["alavancagem_x"] == indicators["alavancagem_x"] and indicators["alavancagem_x"] > dr.THRESHOLDS["alavancagem_max_x"]:
        alertas.append(("Alavancagem acima da referência",
                         f"Dívida Líquida/EBITDA anualizado em {dr.fmt_x(indicators['alavancagem_x'])} — referência usual: até {dr.THRESHOLDS['alavancagem_max_x']:.1f}x."))
    if last["endividamento_geral_pct"] > dr.THRESHOLDS["endividamento_geral_max_pct"]:
        alertas.append(("Endividamento geral elevado", f"Passivo total representa {dr.fmt_pct(last['endividamento_geral_pct'])} do Ativo Total."))
    if last["concentracao_divida_cp_pct"] > dr.THRESHOLDS["concentracao_divida_cp_max_pct"]:
        alertas.append(("Concentração de dívida no curto prazo", f"{dr.fmt_pct(last['concentracao_divida_cp_pct'])} do passivo vence no curto prazo."))
    if last["caixa_operacional_mi"] < 0:
        alertas.append(("Caixa operacional negativo",
                         f"A operação consumiu {dr.fmt_mi(abs(last['caixa_operacional_mi']))} de caixa no último período — investigue a causa antes de apresentar."))

    if not alertas:
        add_rect(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.2), fill=POS_BG)
        add_text(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.7),
                  "Nenhum indicador ultrapassou os limiares de referência configurados no último período coberto.",
                  size=14, color=POS)
    else:
        y = Inches(1.7)
        card_h = Inches(0.85)
        for titulo, texto in alertas[:6]:
            add_rect(s, Inches(0.6), y, Inches(12.1), card_h, fill=CRIT_BG)
            add_text(s, Inches(0.9), y + Inches(0.08), Inches(11.5), Inches(0.35), titulo, size=13, bold=True, color=CRIT)
            add_text(s, Inches(0.9), y + Inches(0.42), Inches(11.5), Inches(0.4), texto, size=11, color=RGBColor(0x3A, 0x20, 0x18))
            y += card_h + Inches(0.12)
    footer(s)

    # ---------------- SLIDE 8: NOTA METODOLÓGICA + ENCERRAMENTO ----------------
    s = blank_slide(prs, bg=NAVY_DARK)
    add_rect(s, Inches(5.87), Inches(0.7), Inches(1.6), Inches(1.6), fill=NAVY)
    add_text(s, Inches(5.87), Inches(0.7), Inches(1.6), Inches(1.6), "i", size=44, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.8), Inches(2.6), Inches(11.73), Inches(0.5), "Nota Metodológica", size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    nota = (f"Consolidação por soma simples (sem eliminação de saldos intercompanhias) das empresas "
            f"informadas, com base nos balancetes de verificação enviados. Indicadores anualizados "
            f"(ROE, ROA, Dívida Líquida/EBITDA) usam fator de anualização de {indicators['fator_anualizacao']:.2f}x "
            f"sobre o período coberto — aproximação simples que não considera sazonalidade. Relatório gerado "
            f"automaticamente; recomenda-se revisão humana antes do envio ao banco.")
    add_text(s, Inches(1.8), Inches(3.3), Inches(9.73), Inches(1.8), nota, size=13, color=RGBColor(0xAF, 0xC1, 0xE8), align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(5.6), Inches(11.73), Inches(0.4), f"Elaborado com apoio de {prepared_by}.",
              size=12, italic=True, color=RGBColor(0x8F, 0xA5, 0xD6), align=PP_ALIGN.CENTER)

    prs.save(out_path)
    return out_path
